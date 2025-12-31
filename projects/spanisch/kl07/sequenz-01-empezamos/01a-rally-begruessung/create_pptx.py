#!/usr/bin/env python3
"""
PowerPoint-Generator für LP-01a: Begrüßung auf Spanisch
Greenhouse Schools Rostock - Spanisch Klasse 7

Format: 16:9 Widescreen
Design: Minimalistisch, CLAUDE.md-konform
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ============================================================================
# DESIGN-KONSTANTEN
# ============================================================================

# Farben (CLAUDE.md-konform)
C_SPANISCH = RGBColor(0xc4, 0x1e, 0x3a)  # Karminrot
C_TEXT = RGBColor(0x22, 0x22, 0x22)       # Fast-Schwarz
C_GRAU = RGBColor(0x66, 0x66, 0x66)       # Mittelgrau
C_HELLGRAU = RGBColor(0xf5, 0xf5, 0xf5)   # Hintergrund
C_WEISS = RGBColor(0xff, 0xff, 0xff)
C_GRUEN = RGBColor(0x2a, 0x7a, 0x4b)      # Erfolg
C_ORANGE = RGBColor(0xb3, 0x5c, 0x00)     # Warnung/Hinweis

# 16:9 Dimensionen
SLIDE_WIDTH = Cm(33.867)   # 16:9 Standard
SLIDE_HEIGHT = Cm(19.05)

# Layout-Konstanten
MARGIN_LEFT = Cm(1.5)
MARGIN_TOP = Cm(1.2)
CONTENT_WIDTH = Cm(30.8)
FOOTER_TOP = Cm(17.8)

# Schriftgrößen
FONT_TITLE = Pt(44)
FONT_SUBTITLE = Pt(28)
FONT_BODY = Pt(22)
FONT_SMALL = Pt(16)
FONT_FOOTER = Pt(11)

# ============================================================================
# HILFSFUNKTIONEN
# ============================================================================

def create_presentation():
    """Erstellt 16:9 Präsentation"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_slide(prs):
    """Fügt leere Folie hinzu"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Weißer Hintergrund
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WEISS
    return slide

def add_title(slide, text, top=MARGIN_TOP, size=FONT_TITLE, color=C_SPANISCH):
    """Haupttitel"""
    shape = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Cm(1.8))
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Helvetica Neue"
    return shape

def add_subtitle(slide, text, top=Cm(3.2), size=FONT_SUBTITLE, color=C_TEXT):
    """Untertitel"""
    shape = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Cm(1.5))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.name = "Helvetica Neue"
    return shape

def add_body_text(slide, text, left, top, width, height, size=FONT_BODY,
                  color=C_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Fließtext"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica Neue"
    p.alignment = align
    return shape

def add_vocab_row(slide, spanish, german, extra, top, highlight=False):
    """Einzelne Vokabelzeile mit Icons"""
    # Spanisch (fett, rot)
    add_body_text(slide, spanish, Cm(2.5), top, Cm(9), Cm(1.2),
                  size=Pt(26), color=C_SPANISCH, bold=True)
    # Deutsch
    add_body_text(slide, german, Cm(12), top, Cm(8), Cm(1.2),
                  size=Pt(24), color=C_TEXT)
    # Extra (Tageszeit etc.)
    if extra:
        add_body_text(slide, extra, Cm(21), top, Cm(11), Cm(1.2),
                      size=Pt(20), color=C_GRAU)

def add_table(slide, data, headers, top, col_widths=None):
    """Erstellt formatierte Tabelle"""
    rows = len(data) + 1
    cols = len(headers)

    if col_widths is None:
        col_widths = [Cm(10)] * cols

    total_width = sum(col_widths)
    left = (SLIDE_WIDTH - total_width) / 2  # Zentriert
    height = Cm(1.1 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, total_width, height).table

    # Spaltenbreiten setzen
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_SPANISCH
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_WEISS
        p.font.name = "Helvetica Neue"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Daten
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.name = "Helvetica Neue"
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Erste Spalte: Spanisch in Rot
            if col_idx == 0:
                p.font.color.rgb = C_SPANISCH
                p.font.bold = True
            else:
                p.font.color.rgb = C_TEXT

            # Zebra-Streifen
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_HELLGRAU

    return table

def add_info_box(slide, title, content, top, icon="💡"):
    """Info-Box mit farbigem Rahmen"""
    left = Cm(2)
    width = Cm(29.8)
    height = Cm(3.5)

    # Hintergrund
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xff, 0xf8, 0xf0)  # Warm-Weiß
    rect.line.color.rgb = C_ORANGE
    rect.line.width = Pt(2)

    # Icon + Titel
    add_body_text(slide, f"{icon} {title}", left + Cm(0.5), top + Cm(0.3),
                  width - Cm(1), Cm(1), size=Pt(18), color=C_ORANGE, bold=True)

    # Inhalt
    add_body_text(slide, content, left + Cm(0.5), top + Cm(1.3),
                  width - Cm(1), height - Cm(1.5), size=Pt(18), color=C_TEXT)

def add_ab_reference(slide, text, top):
    """AB-Referenz-Box"""
    left = Cm(2)
    width = Cm(29.8)
    height = Cm(1.8)

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xfd, 0xf2, 0xf4)
    rect.line.color.rgb = C_SPANISCH
    rect.line.width = Pt(1.5)

    add_body_text(slide, f"📝 {text}", left + Cm(0.5), top + Cm(0.35),
                  width - Cm(1), Cm(1.2), size=Pt(16), color=C_SPANISCH, bold=True)

def add_footer(slide, left_text="Spanisch Klasse 7", right_text="Greenhouse Schools"):
    """Fußzeile mit zwei Spalten"""
    # Links
    add_body_text(slide, left_text, MARGIN_LEFT, FOOTER_TOP, Cm(12), Cm(0.8),
                  size=FONT_FOOTER, color=C_GRAU)
    # Rechts
    add_body_text(slide, right_text, Cm(20), FOOTER_TOP, Cm(12), Cm(0.8),
                  size=FONT_FOOTER, color=C_GRAU, align=PP_ALIGN.RIGHT)
    # Linie
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_LEFT, FOOTER_TOP - Cm(0.15),
                                   CONTENT_WIDTH, Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = C_HELLGRAU
    line.line.fill.background()

def add_map_placeholder(slide, top=Cm(5)):
    """Platzhalter für Weltkarte (manuell einfügen)"""
    left = Cm(4)
    width = Cm(25.8)
    height = Cm(10)

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_HELLGRAU
    rect.line.color.rgb = C_GRAU
    rect.line.width = Pt(1)
    rect.line.dash_style = 2  # Dashed

    # Hinweistext
    add_body_text(slide, "🗺️ Hier Karte einfügen:", left + Cm(0.5), top + Cm(3.5),
                  width - Cm(1), Cm(1.5), size=Pt(20), color=C_GRAU, align=PP_ALIGN.CENTER)
    add_body_text(slide, "Map-Hispanophone_World.svg", left + Cm(0.5), top + Cm(5),
                  width - Cm(1), Cm(1), size=Pt(16), color=C_GRAU, align=PP_ALIGN.CENTER)
    add_body_text(slide, "(Download: commons.wikimedia.org)", left + Cm(0.5), top + Cm(6),
                  width - Cm(1), Cm(1), size=Pt(14), color=C_GRAU, align=PP_ALIGN.CENTER)

# ============================================================================
# FOLIEN
# ============================================================================

def slide_01_title(prs):
    """Folie 1: Titelfolie"""
    slide = add_slide(prs)

    # Großer Titel
    add_body_text(slide, "¡Hola!", Cm(2), Cm(5), Cm(30), Cm(3),
                  size=Pt(96), color=C_SPANISCH, bold=True, align=PP_ALIGN.CENTER)

    # Untertitel
    add_body_text(slide, "Begrüßung und Verabschiedung auf Spanisch",
                  Cm(2), Cm(9.5), Cm(30), Cm(1.5),
                  size=Pt(32), color=C_TEXT, align=PP_ALIGN.CENTER)

    # Metadaten
    add_body_text(slide, "Sequenz 1 · Stunde 1a · Klasse 7",
                  Cm(2), Cm(11.5), Cm(30), Cm(1),
                  size=Pt(20), color=C_GRAU, align=PP_ALIGN.CENTER)

    # Flaggen
    add_body_text(slide, "🇪🇸   🇲🇽   🇦🇷   🇨🇴   🇵🇪   🇨🇺   🇨🇱   🇻🇪",
                  Cm(2), Cm(14), Cm(30), Cm(1.5),
                  size=Pt(36), color=C_TEXT, align=PP_ALIGN.CENTER)

    add_footer(slide)

def slide_02_map(prs):
    """Folie 2: Weltkarte"""
    slide = add_slide(prs)

    add_title(slide, "🌍 Wo spricht man Spanisch?")
    add_subtitle(slide, "Spanisch ist Amtssprache in 21 Ländern – gesprochen von über 480 Mio. Menschen!")

    # Karten-Platzhalter
    add_map_placeholder(slide, top=Cm(5))

    add_footer(slide)

def slide_03_greeting(prs):
    """Folie 3: Begrüßungen"""
    slide = add_slide(prs)

    add_title(slide, "☀️ Saludar – Begrüßen")
    add_subtitle(slide, "Je nach Tageszeit verwendet man unterschiedliche Begrüßungen:")

    # Tabelle
    data = [
        ("🙋 ¡Hola!", "Hallo!", "Immer (informell)"),
        ("🌅 Buenos días", "Guten Morgen", "bis ca. 14 Uhr"),
        ("🌇 Buenas tardes", "Guten Tag", "14–20 Uhr"),
        ("🌙 Buenas noches", "Guten Abend/Nacht", "ab 20 Uhr"),
    ]
    add_table(slide, data, ["Spanisch", "Deutsch", "Wann?"], Cm(5),
              col_widths=[Cm(11), Cm(9), Cm(8)])

    # AB-Referenz
    add_ab_reference(slide, "AB-01a: Kasten 1 – Begrüßungen eintragen", Cm(14))

    add_footer(slide)

def slide_04_farewell(prs):
    """Folie 4: Verabschiedungen"""
    slide = add_slide(prs)

    add_title(slide, "👋 Despedirse – Verabschieden")
    add_subtitle(slide, "So verabschiedet man sich auf Spanisch:")

    data = [
        ("¡Adiós!", "Tschüss! / Auf Wiedersehen!"),
        ("¡Hasta luego!", "Bis später!"),
        ("¡Hasta mañana!", "Bis morgen!"),
        ("¡Hasta pronto!", "Bis bald!"),
    ]
    add_table(slide, data, ["Spanisch", "Deutsch"], Cm(5),
              col_widths=[Cm(12), Cm(14)])

    # AB-Referenz
    add_ab_reference(slide, "AB-01a: Kasten 1 – Verabschiedungen ergänzen", Cm(13))

    add_footer(slide)

def slide_05_grammar(prs):
    """Folie 5: Grammatik"""
    slide = add_slide(prs)

    add_title(slide, "📝 Grammatik-Tipp")

    # buenos vs buenas
    add_body_text(slide, "Buenos días", Cm(3), Cm(4.5), Cm(12), Cm(1.5),
                  size=Pt(40), color=C_SPANISCH, bold=True)
    add_body_text(slide, "← maskulin (el día)", Cm(16), Cm(4.8), Cm(14), Cm(1),
                  size=Pt(24), color=C_GRAU)

    add_body_text(slide, "Buenas tardes / noches", Cm(3), Cm(7), Cm(16), Cm(1.5),
                  size=Pt(40), color=C_SPANISCH, bold=True)
    add_body_text(slide, "← feminin (la tarde, la noche)", Cm(20), Cm(7.3), Cm(12), Cm(1),
                  size=Pt(24), color=C_GRAU)

    # Info-Box
    add_info_box(slide, "Merke",
                 "Das Adjektiv passt sich dem Geschlecht des Substantivs an:\nbueno → buenos (m.)  /  buena → buenas (f.)",
                 Cm(10))

    # Satzzeichen
    add_body_text(slide, "¡Hola!  ¿Qué tal?", Cm(3), Cm(15), Cm(20), Cm(1.2),
                  size=Pt(36), color=C_SPANISCH, bold=True)
    add_body_text(slide, "→  Satzzeichen am Anfang UND Ende!", Cm(3), Cm(16.5), Cm(25), Cm(1),
                  size=Pt(20), color=C_GRAU)

    add_footer(slide)

def slide_06_exercise(prs):
    """Folie 6: Übung Zuordnung"""
    slide = add_slide(prs)

    add_title(slide, "✏️ Übung: Zuordnung")
    add_subtitle(slide, "Ordne die spanischen Ausdrücke der deutschen Übersetzung zu:")

    # Übung
    exercises = [
        ("1.  ¡Hola!", "____"),
        ("2.  Buenos días", "____"),
        ("3.  ¡Hasta mañana!", "____"),
        ("4.  Buenas noches", "____"),
    ]

    for i, (spanish, blank) in enumerate(exercises):
        top = Cm(5.5 + i * 2.2)
        add_body_text(slide, spanish, Cm(4), top, Cm(11), Cm(1.2),
                      size=Pt(28), color=C_SPANISCH, bold=True)
        add_body_text(slide, blank, Cm(17), top, Cm(8), Cm(1.2),
                      size=Pt(28), color=C_GRAU)

    # Optionen
    add_body_text(slide, "Optionen:   a) Guten Abend   b) Bis morgen!   c) Guten Morgen   d) Hallo!",
                  Cm(4), Cm(14.5), Cm(26), Cm(1), size=Pt(18), color=C_GRAU)

    # AB-Referenz
    add_ab_reference(slide, "AB-01a: Aufgabe 1 – Zuordnung (4 Punkte)", Cm(16))

    add_footer(slide)

def slide_07_solution(prs):
    """Folie 7: Lösung"""
    slide = add_slide(prs)

    add_title(slide, "✅ Lösung")

    solutions = [
        ("1.  ¡Hola!", "→  d) Hallo!"),
        ("2.  Buenos días", "→  c) Guten Morgen"),
        ("3.  ¡Hasta mañana!", "→  b) Bis morgen!"),
        ("4.  Buenas noches", "→  a) Guten Abend"),
    ]

    for i, (spanish, solution) in enumerate(solutions):
        top = Cm(4.5 + i * 2.5)
        add_body_text(slide, spanish, Cm(4), top, Cm(11), Cm(1.2),
                      size=Pt(28), color=C_SPANISCH, bold=True)
        add_body_text(slide, solution, Cm(16), top, Cm(14), Cm(1.2),
                      size=Pt(28), color=C_GRUEN, bold=True)

    add_footer(slide)

def slide_08_exercise2(prs):
    """Folie 8: Übung Tageszeiten"""
    slide = add_slide(prs)

    add_title(slide, "🕐 Übung: Tageszeiten")
    add_subtitle(slide, "Welche Begrüßung passt? Schreibe das richtige spanische Wort:")

    times = [
        ("a)  8:00 Uhr morgens", "________________"),
        ("b)  15:00 Uhr nachmittags", "________________"),
        ("c)  21:00 Uhr abends", "________________"),
        ("d)  Jederzeit (informell)", "________________"),
    ]

    for i, (time, blank) in enumerate(times):
        top = Cm(5.5 + i * 2.2)
        add_body_text(slide, time, Cm(4), top, Cm(14), Cm(1.2),
                      size=Pt(26), color=C_TEXT)
        add_body_text(slide, blank, Cm(19), top, Cm(12), Cm(1.2),
                      size=Pt(26), color=C_GRAU)

    # AB-Referenz
    add_ab_reference(slide, "AB-01a: Aufgabe 2 – Tageszeiten (4 Punkte)", Cm(15))

    add_footer(slide)

def slide_09_summary(prs):
    """Folie 9: Zusammenfassung"""
    slide = add_slide(prs)

    add_title(slide, "📋 Zusammenfassung")

    # Zwei Spalten
    add_body_text(slide, "Begrüßung (Saludar)", Cm(2.5), Cm(4), Cm(14), Cm(1.2),
                  size=Pt(24), color=C_SPANISCH, bold=True)

    greetings = ["¡Hola!", "Buenos días", "Buenas tardes", "Buenas noches"]
    for i, g in enumerate(greetings):
        add_body_text(slide, f"•  {g}", Cm(3), Cm(5.5 + i * 1.6), Cm(12), Cm(1.2),
                      size=Pt(22), color=C_TEXT)

    add_body_text(slide, "Verabschiedung (Despedirse)", Cm(17.5), Cm(4), Cm(14), Cm(1.2),
                  size=Pt(24), color=C_SPANISCH, bold=True)

    farewells = ["¡Adiós!", "¡Hasta luego!", "¡Hasta mañana!", "¡Hasta pronto!"]
    for i, f in enumerate(farewells):
        add_body_text(slide, f"•  {f}", Cm(18), Cm(5.5 + i * 1.6), Cm(12), Cm(1.2),
                      size=Pt(22), color=C_TEXT)

    # Merksatz
    add_info_box(slide, "Denke daran",
                 "buenos (m.) vs. buenas (f.)  •  Satzzeichen: ¡...!  ¿...?  •  21 Länder sprechen Spanisch!",
                 Cm(12.5))

    add_footer(slide)

def slide_10_homework(prs):
    """Folie 10: Hausaufgabe"""
    slide = add_slide(prs)

    add_title(slide, "📚 Hausaufgabe")

    tasks = [
        ("✓", "AB-01a fertig ausfüllen (Aufgaben 3 + 4)"),
        ("✓", "Lehrwerk: Qué pasa 1, S. 10–15 lesen"),
        ("✓", "Vokabeln lernen: Begrüßung + Verabschiedung"),
        ("✓", "Mini-Dialog mit Partner/in üben (Aufgabe 4)"),
    ]

    for i, (check, task) in enumerate(tasks):
        top = Cm(4.5 + i * 2.3)
        add_body_text(slide, check, Cm(3), top, Cm(1.5), Cm(1.2),
                      size=Pt(28), color=C_GRUEN, bold=True)
        add_body_text(slide, task, Cm(5), top, Cm(26), Cm(1.2),
                      size=Pt(24), color=C_TEXT)

    # AB-Referenz
    add_ab_reference(slide, "AB-01a: Aufgabe 3 – Rally durchs Buch (6P) · Aufgabe 4 – Mini-Dialog (6P)", Cm(14))

    # Abschluss
    add_body_text(slide, "¡Hasta la próxima clase!", Cm(2), Cm(16.5), Cm(30), Cm(1.5),
                  size=Pt(36), color=C_SPANISCH, bold=True, align=PP_ALIGN.CENTER)

    add_footer(slide)

# ============================================================================
# HAUPTPROGRAMM
# ============================================================================

def main():
    """Erstellt die Präsentation"""
    print("Erstelle PowerPoint-Präsentation (16:9)...")

    prs = create_presentation()

    # Alle Folien erstellen
    slide_01_title(prs)
    slide_02_map(prs)
    slide_03_greeting(prs)
    slide_04_farewell(prs)
    slide_05_grammar(prs)
    slide_06_exercise(prs)
    slide_07_solution(prs)
    slide_08_exercise2(prs)
    slide_09_summary(prs)
    slide_10_homework(prs)

    # Speichern
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "PPT-01a-begruessung.pptx")
    prs.save(output_path)

    print(f"✓ Gespeichert: {output_path}")
    print(f"  → 10 Folien im 16:9-Format")
    print()
    print("HINWEIS: Folie 2 enthält Platzhalter für Weltkarte.")
    print("         Download: https://commons.wikimedia.org/wiki/File:Map-Hispanophone_World.svg")

if __name__ == "__main__":
    main()
