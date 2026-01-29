#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Generator: Repaso Unidad 5 + Números 100-1000
Spanisch Klasse 8 - 90-Minuten-Doppelstunde
Schüler-gerichtete Präsentation mit Aufgaben + Lösungen auf Folgefolien
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Farben
SPANISCH_ROT = RGBColor(0xC4, 0x1E, 0x3A)
LOESUNG_ROT = RGBColor(0xC6, 0x28, 0x28)
DUNKELGRAU = RGBColor(0x33, 0x33, 0x33)
HELLGRAU = RGBColor(0xF5, 0xF5, 0xF5)
WEISS = RGBColor(0xFF, 0xFF, 0xFF)
GRUEN = RGBColor(0x2A, 0x7A, 0x4B)

def add_header(slide, title, subtitle=""):
    """Fügt Header-Balken mit Titel hinzu."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SPANISCH_ROT
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WEISS

    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(9), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = WEISS

def add_content_lines(slide, lines, start_top=1.4, color=DUNKELGRAU):
    """Fügt Textzeilen hinzu."""
    top = start_top
    for line in lines:
        if line == "---":
            top += 0.15
            continue

        is_bold = line.startswith("**")
        is_bullet = line.startswith("•")
        is_spanish = line.startswith("[ES]")

        text = line.replace("**", "").replace("[ES]", "").strip()
        left = 0.7 if is_bullet else 0.5

        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(9), Inches(0.45))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20) if is_bold else Pt(18)
        p.font.bold = is_bold
        p.font.color.rgb = SPANISCH_ROT if (is_bold or is_spanish) else color
        top += 0.45
    return top

def add_solution_header(slide):
    """Fügt Lösungs-Badge hinzu."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(0.15), Inches(2), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRUEN
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(7.8), Inches(0.22), Inches(2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "LÖSUNG"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WEISS
    p.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ===== FOLIE 1: TITEL =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SPANISCH_ROT
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Repaso Unidad 5"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = SPANISCH_ROT
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "La ropa, los colores, los números 100-1000"
    p.font.size = Pt(24)
    p.font.color.rgb = DUNKELGRAU
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Spanisch Klasse 8"
    p.font.size = Pt(18)
    p.font.color.rgb = DUNKELGRAU
    p.alignment = PP_ALIGN.CENTER

    # ===== FOLIE 2: EINSTIEG =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "¡Vamos! Schätz-Spiel")
    add_content_lines(slide, [
        "**¿Cuánto cuesta? – Was kostet das?**",
        "---",
        "Schau dir die Gegenstände an und schätze den Preis!",
        "Antworte auf Spanisch.",
        "---",
        "• Un libro (ein Buch)",
        "• Una camiseta (ein T-Shirt)",
        "• Unos zapatos (Schuhe)",
        "---",
        '[ES]Ejemplo: "Creo que cuesta... veinte euros."'
    ])

    # ===== FOLIE 3: KLEIDUNG INPUT =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "La ropa – Kleidung", "Buch S. 78-79")
    add_content_lines(slide, [
        "**Diese Vokabeln brauchst du:**",
        "---",
        "[ES]la camiseta = das T-Shirt",
        "[ES]el jersey = der Pullover",
        "[ES]los pantalones = die Hose",
        "[ES]el vestido = das Kleid",
        "[ES]los zapatos = die Schuhe",
        "[ES]la falda = der Rock"
    ])

    # ===== FOLIE 4: AUFGABE 1 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 1: La ropa", "Arbeitsblatt • 5 Punkte")
    add_content_lines(slide, [
        "**Verbinde mit Linien!**",
        "---",
        "Zeichne eine Linie vom deutschen zum spanischen Wort.",
        "---",
        "Du hast 2 Minuten Zeit.",
        "---",
        "[ES]¡Vamos!"
    ])

    # ===== FOLIE 5: LÖSUNG 1 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 1: La ropa", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "**Die richtigen Paare:**",
        "---",
        "[ES]T-Shirt → la camiseta",
        "[ES]Hose → los pantalones",
        "[ES]Schuhe → los zapatos",
        "[ES]Kleid → el vestido",
        "[ES]Pullover → el jersey"
    ], color=GRUEN)

    # ===== FOLIE 6: FARBEN INPUT =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Los colores – Farben", "Achtung: Endungen!")
    add_content_lines(slide, [
        "**Farben passen sich dem Nomen an!**",
        "---",
        "[ES]blanco / blanca = weiß",
        "[ES]negro / negra = schwarz",
        "[ES]rojo / roja = rot",
        "[ES]azul = blau (keine Änderung)",
        "[ES]marrón = braun (keine Änderung)",
        "---",
        "Beispiel: la camiseta blanc**a** (feminin!)"
    ])

    # ===== FOLIE 7: AUFGABE 2 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 2: Los colores", "Arbeitsblatt • 4 Punkte")
    add_content_lines(slide, [
        "**Schreibe die richtige Farbe!**",
        "---",
        "Achte auf die Endung: -o oder -a?",
        "---",
        "Schau auf den Artikel: el = maskulin, la = feminin",
        "---",
        "[ES]¡Adelante!"
    ])

    # ===== FOLIE 8: LÖSUNG 2 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 2: Los colores", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "**Die richtigen Farben:**",
        "---",
        "[ES]a) la camiseta blanca (weiß, feminin)",
        "[ES]b) el jersey negro (schwarz, maskulin)",
        "[ES]c) los zapatos marrón (braun, invariabel)",
        "[ES]d) la falda roja (rot, feminin)"
    ], color=GRUEN)

    # ===== FOLIE 9: MODALVERBEN INPUT =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Modalverben", "Buch S. 82")
    add_content_lines(slide, [
        "**querer (wollen) – poder (können) – tener que (müssen)**",
        "---",
        "[ES]yo quiero / puedo / tengo que",
        "[ES]tú quieres / puedes / tienes que",
        "[ES]él/ella quiere / puede / tiene que",
        "[ES]nosotros queremos / podemos / tenemos que",
        "[ES]ellos quieren / pueden / tienen que"
    ])

    # ===== FOLIE 10: AUFGABE 3 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 3: Modalverben", "Arbeitsblatt • 5 Punkte")
    add_content_lines(slide, [
        "**Ergänze die richtige Form!**",
        "---",
        "Schau auf das Subjekt: yo, tú, nosotros...?",
        "Welches Verb steht in Klammern?",
        "---",
        "Du hast 3 Minuten Zeit.",
        "---",
        "[ES]¡Mucha suerte!"
    ])

    # ===== FOLIE 11: LÖSUNG 3 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 3: Modalverben", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "[ES]a) Yo quiero comprar un jersey.",
        "[ES]b) Tú no puedes ir a la fiesta.",
        "[ES]c) Nosotros tenemos que estudiar.",
        "[ES]d) Ellos quieren una chaqueta.",
        "[ES]e) María puede hablar español."
    ], color=GRUEN)

    # ===== FOLIE 12: ZAHLEN 100-1000 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Los números 100-1000", "Neue Zahlen!")

    # Merkkasten
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(9.4), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HELLGRAU
    shape.line.color.rgb = SPANISCH_ROT
    shape.line.width = Pt(2)

    add_content_lines(slide, [
        "[ES]100 = cien     200 = doscientos/as     300 = trescientos/as",
        "[ES]400 = cuatrocientos/as     500 = quinientos/as",
        "[ES]600 = seiscientos/as     700 = setecientos/as",
        "[ES]800 = ochocientos/as     900 = novecientos/as     1000 = mil",
        "---",
        "**Wichtig:** 100 allein = cien, aber 101 = ciento uno",
        "**Unregelmäßig:** quinientos, setecientos, novecientos",
        "**Genus:** doscientos euros / doscientas personas"
    ], start_top=1.5)

    # ===== FOLIE 13: AUFGABE 4 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 4: Zahlen schreiben", "Arbeitsblatt • 6 Punkte")
    add_content_lines(slide, [
        "**Schreibe die Zahlen als spanische Wörter!**",
        "---",
        "Nutze den Merkkasten auf deinem Arbeitsblatt.",
        "---",
        "Tipp: Große Zahlen = Hunderter + Zehner/Einer",
        "Beispiel: 725 = setecientos + veinticinco",
        "---",
        "[ES]¡Inténtalo!"
    ])

    # ===== FOLIE 14: LÖSUNG 4 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 4: Zahlen schreiben", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "[ES]a) 100 = cien",
        "[ES]b) 350 = trescientos cincuenta",
        "[ES]c) 500 = quinientos",
        "[ES]d) 725 = setecientos veinticinco",
        "[ES]e) 999 = novecientos noventa y nueve",
        "[ES]f) 1000 = mil"
    ], color=GRUEN)

    # ===== FOLIE 15: AUFGABE 5 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 5: Zahlen erkennen", "Arbeitsblatt • 4 Punkte")
    add_content_lines(slide, [
        "**Schreibe die Zahl als Ziffer!**",
        "---",
        "Lies das spanische Wort und schreibe die Zahl.",
        "---",
        "Tipp: Zerlege in Hunderter + Rest",
        "---",
        "[ES]¡Es fácil!"
    ])

    # ===== FOLIE 16: LÖSUNG 5 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 5: Zahlen erkennen", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "[ES]a) doscientos cuarenta = 240",
        "[ES]b) quinientos ochenta y tres = 583",
        "[ES]c) setecientos quince = 715",
        "[ES]d) novecientos noventa y nueve = 999"
    ], color=GRUEN)

    # ===== FOLIE 17: EINKAUFSDIALOG INPUT =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "De compras – Einkaufen", "Buch S. 84-85")
    add_content_lines(slide, [
        "**Wichtige Redemittel:**",
        "---",
        "[ES]¿En qué puedo ayudarte? = Wie kann ich dir helfen?",
        "[ES]Quiero comprar... = Ich möchte ... kaufen",
        "[ES]¿Qué color quieres? = Welche Farbe möchtest du?",
        "[ES]¿Cuánto cuesta? = Wieviel kostet das?",
        "[ES]Lo compro. = Ich kaufe es."
    ])

    # ===== FOLIE 18: AUFGABE 6 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 6: Einkaufsdialog", "Arbeitsblatt • 4 Punkte")
    add_content_lines(slide, [
        "**Ergänze den Dialog auf Spanisch!**",
        "---",
        "Du bist der Kunde (Cliente) und möchtest einkaufen.",
        "Die deutschen Übersetzungen helfen dir.",
        "---",
        "Nutze die Redemittel von der vorherigen Folie!",
        "---",
        "[ES]¡Buena suerte!"
    ])

    # ===== FOLIE 19: LÖSUNG 6 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Aufgabe 6: Einkaufsdialog", "Lösung")
    add_solution_header(slide)
    add_content_lines(slide, [
        "[ES]Cliente: Quiero comprar un jersey.",
        "[ES]Cliente: Quiero uno azul.",
        "[ES]Cliente: ¿Cuánto cuesta?",
        "[ES]Cliente: Bien, lo compro."
    ], color=GRUEN)

    # ===== FOLIE 20: EXTENSION TASKS =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Extension Tasks (★★★)", "Für Schnelle • +5 BE")
    add_content_lines(slide, [
        "**E1: Rechenaufgabe (+2 BE)**",
        "Pablo kauft eine Jacke für 245€ und Schuhe für 189€.",
        "Schreibe den Gesamtpreis als spanisches Wort.",
        "---",
        "**E2: Eigener Dialog (+3 BE)**",
        "Schreibe einen Einkaufsdialog (4 Zeilen):",
        "Du möchtest ein rotes Kleid kaufen, das 599€ kostet.",
        "---",
        "[ES]¡Buena suerte!"
    ])

    # ===== FOLIE 21: EXTENSION LÖSUNGEN =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Extension Tasks (★★★)", "Lösungen")
    add_solution_header(slide)
    add_content_lines(slide, [
        "**E1:** 245 + 189 = 434",
        "[ES]cuatrocientos treinta y cuatro",
        "---",
        "**E2:** Beispieldialog:",
        "[ES]Cliente: Quiero comprar un vestido rojo.",
        "[ES]Vendedor: Aquí tienes. Cuesta 599 euros.",
        "[ES]Cliente: ¿Cuánto cuesta?",
        "[ES]Vendedor: Quinientos noventa y nueve euros."
    ], color=GRUEN)

    # ===== FOLIE 22: ABSCHLUSS =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "¡Muy bien!", "Das hast du heute gelernt:")
    add_content_lines(slide, [
        "**Vokabeln:**",
        "• La ropa (Kleidung) + Los colores (Farben)",
        "• Los números 100-1000",
        "---",
        "**Grammatik:**",
        "• Modalverben: querer, poder, tener que",
        "• Farben mit Genus-Anpassung",
        "---",
        "**Kommunikation:**",
        "• Einkaufsdialog führen"
    ])

    # ===== FOLIE 23: BEWERTUNG (für Lehrkraft) =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Bewertung", "28 BE + 5 BE Extension")
    add_content_lines(slide, [
        "**Basis (28 BE):**",
        "• Aufgabe 1-6: 5+4+5+6+4+4 = 28 BE",
        "---",
        "**Extension (★★★, +5 BE):**",
        "• E1 Rechenaufgabe: +2 BE",
        "• E2 Eigener Dialog: +3 BE",
        "---",
        "**14-NP-Skala:** 28 BE = 14 NP | 22 BE = 10 NP | 13 BE = 5 NP"
    ])

    # Speichern
    prs.save('PRAESENTATION-repaso-numeros.pptx')
    print("✓ PPTX erstellt: PRAESENTATION-repaso-numeros.pptx (23 Folien)")
    print("  - Schülergerichtete Ansprache")
    print("  - Aufgaben und Lösungen auf getrennten Folien")
    print("  - Extension Tasks für Gymnasialstufe")

if __name__ == "__main__":
    main()
