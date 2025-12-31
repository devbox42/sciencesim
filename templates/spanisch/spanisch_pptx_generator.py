#!/usr/bin/env python3
"""
Spanisch PowerPoint Generator für Klasse 7
Greenhouse Schools Rostock

Basiert auf ChalkSheet Template - getestet ohne Repair-Warnungen.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Farben
PRIMARY = RGBColor(0xc4, 0x1e, 0x3a)      # Karminrot
PRIMARY_LIGHT = RGBColor(0xfd, 0xf2, 0xf4)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
MEDIUM_TEXT = RGBColor(0x55, 0x55, 0x55)
LIGHT_TEXT = RGBColor(0x88, 0x88, 0x88)
SUCCESS = RGBColor(0x2a, 0x7a, 0x4b)
WARNING = RGBColor(0xb3, 0x5c, 0x00)


class SpanischPPTX:
    def __init__(self, lesson_id, title, subtitle="", sequenz=1):
        self.lesson_id = lesson_id
        self.title = title
        self.subtitle = subtitle
        self.sequenz = sequenz

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _text(self, slide, text, left, top, width, height, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        if color is None:
            color = DARK_TEXT
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        para = box.text_frame.paragraphs[0]
        para.text = text
        para.font.name = "Arial"
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = color
        para.alignment = align
        return box

    def _rect(self, slide, left, top, width, height, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line:
            shape.line.color.rgb = line
        else:
            shape.line.fill.background()
        return shape

    def _header_footer(self, slide):
        # Header
        self._rect(slide, 0, 0, 13.333, 0.5, PRIMARY)
        self._text(slide, f"Spanisch Klasse 7 · LP-{self.lesson_id}", 0.3, 0.1, 4, 0.3,
                   size=14, bold=True, color=WHITE)
        # Footer
        self._rect(slide, 0, 7.2, 13.333, 0.3, LIGHT_GRAY)
        self._text(slide, self.title, 0.3, 7.22, 5, 0.25, size=10, color=MEDIUM_TEXT)
        self._text(slide, f"Sequenz {self.sequenz}", 10, 7.22, 3, 0.25,
                   size=10, color=MEDIUM_TEXT, align=PP_ALIGN.RIGHT)

    def _title(self, slide, text, top=0.7):
        self._text(slide, text, 0.5, top, 12.33, 0.8, size=32, bold=True, color=PRIMARY)

    def _badge(self, slide, text, left=0.5, top=0.7, color=None):
        if color is None:
            color = PRIMARY
        self._rect(slide, left, top, 2, 0.5, color)
        self._text(slide, text, left + 0.1, top + 0.08, 1.8, 0.35,
                   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    def _ab_ref(self, slide, text, top=6.5):
        self._rect(slide, 0.5, top, 12.33, 0.6, PRIMARY_LIGHT, PRIMARY)
        self._text(slide, f"📝 {text}", 0.7, top + 0.12, 12, 0.4,
                   size=14, bold=True, color=PRIMARY)

    # === SLIDE TYPES ===

    def add_title_slide(self, flags="🇪🇸"):
        slide = self._slide()
        self._rect(slide, 0, 0, 13.333, 4.5, PRIMARY)
        self._text(slide, self.title, 0.5, 1.5, 12.33, 1.2, size=56, bold=True, color=WHITE)
        if self.subtitle:
            self._text(slide, self.subtitle, 0.5, 2.8, 12.33, 0.6, size=24, color=RGBColor(0xff,0xcc,0xcc))
        self._text(slide, f"Sequenz {self.sequenz} · Lektion {self.lesson_id}", 0.5, 3.6, 12.33, 0.5,
                   size=14, color=RGBColor(0xee,0xaa,0xaa))
        self._text(slide, "Spanisch Klasse 7", 0.5, 5.5, 6, 0.5, size=20, bold=True, color=PRIMARY)
        if flags:
            self._text(slide, flags, 0.5, 6.2, 12.33, 0.8, size=36)
        return slide

    def add_vocab_slide(self, title, vocab, subtitle=None, ab_ref=None):
        slide = self._slide()
        self._header_footer(slide)
        self._title(slide, f"📚 {title}")
        if subtitle:
            self._text(slide, subtitle, 0.5, 1.3, 12.33, 0.5, size=18, color=MEDIUM_TEXT)

        start_y = 1.8 if subtitle else 1.5
        row_h = 0.65

        # Header row
        self._rect(slide, 0.5, start_y, 12.33, row_h, PRIMARY)
        self._text(slide, "Spanisch", 0.7, start_y + 0.12, 5.5, 0.4, size=18, bold=True, color=WHITE)
        self._text(slide, "Deutsch", 6.5, start_y + 0.12, 5.5, 0.4, size=18, bold=True, color=WHITE)

        # Data rows
        for i, item in enumerate(vocab):
            y = start_y + row_h * (i + 1)
            if i % 2 == 1:
                self._rect(slide, 0.5, y, 12.33, row_h, LIGHT_GRAY)
            sp = item[0] if isinstance(item, (list, tuple)) else item
            de = item[1] if isinstance(item, (list, tuple)) and len(item) > 1 else ""
            self._text(slide, sp, 0.7, y + 0.12, 5.5, 0.4, size=20, bold=True, color=PRIMARY)
            self._text(slide, de, 6.5, y + 0.12, 5.5, 0.4, size=20, color=DARK_TEXT)

        if ab_ref:
            self._ab_ref(slide, ab_ref)
        return slide

    def add_grammar_slide(self, title, blocks, info=None):
        slide = self._slide()
        self._header_footer(slide)
        self._title(slide, f"📝 {title}")

        y = 1.6
        for block in blocks:
            if isinstance(block, tuple) and len(block) == 2:
                self._text(slide, block[0], 0.7, y, 5.5, 0.8, size=28, bold=True, color=PRIMARY)
                self._text(slide, block[1], 6.5, y + 0.1, 6, 0.6, size=18, color=MEDIUM_TEXT)
                y += 1.2
            else:
                self._text(slide, str(block), 0.7, y, 12, 0.6, size=20, color=DARK_TEXT)
                y += 0.8

        if info:
            self._rect(slide, 0.5, 5.5, 12.33, 1.2, RGBColor(0xff,0xf8,0xf0), WARNING)
            self._text(slide, f"💡 {info}", 0.7, 5.7, 12, 0.8, size=16, color=DARK_TEXT)
        return slide

    def add_exercise_slide(self, title, exercises, options=None, ab_ref=None):
        slide = self._slide()
        self._header_footer(slide)
        self._badge(slide, "ÜBUNG")
        self._title(slide, f"✏️ {title}", top=1.4)

        y = 2.2
        for left, right in exercises:
            self._text(slide, left, 0.7, y, 5.5, 0.6, size=22, bold=True, color=PRIMARY)
            self._text(slide, right, 6.5, y, 5.5, 0.6, size=22, color=LIGHT_TEXT)
            y += 1.0

        if options:
            self._text(slide, f"Optionen: {options}", 0.7, 5.5, 12, 0.5, size=16, color=MEDIUM_TEXT)
        if ab_ref:
            self._ab_ref(slide, ab_ref)
        return slide

    def add_solution_slide(self, title, solutions):
        slide = self._slide()
        self._header_footer(slide)
        self._badge(slide, "LÖSUNG", color=SUCCESS)
        self._title(slide, f"✅ {title}", top=1.4)

        y = 2.2
        for left, right in solutions:
            self._text(slide, left, 0.7, y, 5.5, 0.6, size=24, bold=True, color=PRIMARY)
            self._text(slide, right, 6.5, y, 5.5, 0.6, size=24, bold=True, color=SUCCESS)
            y += 1.2
        return slide

    def add_summary_slide(self, title, left_items, right_items=None, left_title="Vokabeln", right_title="Grammatik"):
        slide = self._slide()
        self._header_footer(slide)
        self._title(slide, f"📋 {title}")

        self._text(slide, left_title, 0.7, 1.6, 5.5, 0.5, size=20, bold=True, color=PRIMARY)
        y = 2.2
        for item in left_items:
            self._text(slide, f"•  {item}", 0.9, y, 5.3, 0.5, size=18, color=DARK_TEXT)
            y += 0.6

        if right_items:
            self._text(slide, right_title, 7, 1.6, 5.5, 0.5, size=20, bold=True, color=PRIMARY)
            y = 2.2
            for item in right_items:
                self._text(slide, f"•  {item}", 7.2, y, 5.3, 0.5, size=18, color=DARK_TEXT)
                y += 0.6
        return slide

    def add_homework_slide(self, tasks, ab_ref=None, closing="¡Hasta luego!"):
        slide = self._slide()
        self._header_footer(slide)
        self._badge(slide, "HAUSAUFGABE")
        self._title(slide, "📚 Para casa", top=1.4)

        y = 2.3
        for task in tasks:
            self._text(slide, "✓", 0.7, y, 0.5, 0.5, size=24, bold=True, color=SUCCESS)
            self._text(slide, task, 1.3, y, 11, 0.5, size=20, color=DARK_TEXT)
            y += 0.9

        if ab_ref:
            self._ab_ref(slide, ab_ref, top=5.5)
        if closing:
            self._text(slide, closing, 0.5, 6.3, 12.33, 0.6, size=32, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        return slide

    def add_map_slide(self, title, subtitle=None, image_path=None):
        slide = self._slide()
        self._header_footer(slide)
        self._title(slide, f"🌍 {title}")
        if subtitle:
            self._text(slide, subtitle, 0.5, 1.3, 12.33, 0.5, size=18, color=MEDIUM_TEXT)

        if image_path:
            from pathlib import Path
            if Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(1.5), Inches(2), width=Inches(10))
            else:
                self._rect(slide, 1.5, 2, 10, 4.5, LIGHT_GRAY, MEDIUM_TEXT)
                self._text(slide, "[Bild nicht gefunden]", 5, 4, 4, 0.5, size=16, color=MEDIUM_TEXT, align=PP_ALIGN.CENTER)
        return slide

    def add_content_slide(self, title, items, ab_ref=None):
        slide = self._slide()
        self._header_footer(slide)
        self._title(slide, title)

        y = 1.6
        for item in items:
            self._text(slide, item, 0.7, y, 12, 0.6, size=20, color=DARK_TEXT)
            y += 0.8

        if ab_ref:
            self._ab_ref(slide, ab_ref)
        return slide

    def save(self, path):
        self.prs.save(str(path))
        return len(self.prs.slides)


if __name__ == "__main__":
    pptx = SpanischPPTX("TEST", "¡Hola!", "Test-Präsentation", sequenz=1)
    pptx.add_title_slide()
    pptx.add_vocab_slide("Begrüßungen", [("Hola", "Hallo"), ("Adiós", "Tschüss")], ab_ref="AB Kasten 1")
    pptx.add_grammar_slide("Das Verb ser", [("Yo soy", "Ich bin"), ("Tú eres", "Du bist")])
    pptx.add_exercise_slide("Ergänze", [("Buenos ___", "(morgens)"), ("Buenas ___", "(abends)")], options="días, noches")
    pptx.add_solution_slide("Lösung", [("Buenos días", "Guten Morgen"), ("Buenas noches", "Gute Nacht")])
    pptx.add_homework_slide(["Vokabeln lernen", "AB ausfüllen"])
    count = pptx.save("/tmp/test_step5.pptx")
    print(f"✓ Step 5: /tmp/test_step5.pptx ({count} Folien)")
