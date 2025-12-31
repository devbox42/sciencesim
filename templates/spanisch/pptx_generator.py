#!/usr/bin/env python3
"""
Generischer PowerPoint-Generator für Spanisch Klasse 7
Greenhouse Schools Rostock

Verwendung:
    from pptx_generator import PPTXGenerator

    lesson = {
        'id': '01a',
        'title': '¡Hola!',
        'subtitle': 'Begrüßung auf Spanisch',
        'sequenz': 1,
        'block': 'a',
        'slides': [...]  # Folien-Daten
    }

    gen = PPTXGenerator(lesson)
    gen.generate('output.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os


# ============================================================================
# DESIGN-KONSTANTEN (CLAUDE.md-konform)
# ============================================================================

C_SPANISCH = RGBColor(0xc4, 0x1e, 0x3a)    # Karminrot
C_TEXT = RGBColor(0x22, 0x22, 0x22)         # Fast-Schwarz
C_GRAU = RGBColor(0x66, 0x66, 0x66)         # Mittelgrau
C_HELLGRAU = RGBColor(0xf5, 0xf5, 0xf5)     # Hintergrund
C_WEISS = RGBColor(0xff, 0xff, 0xff)
C_GRUEN = RGBColor(0x2a, 0x7a, 0x4b)        # Erfolg
C_ORANGE = RGBColor(0xb3, 0x5c, 0x00)       # Hinweis

# 16:9 Dimensionen
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# Layout
MARGIN_LEFT = Cm(1.5)
MARGIN_TOP = Cm(1.2)
CONTENT_WIDTH = Cm(30.8)
FOOTER_TOP = Cm(17.8)

# Schriftgrößen
FONT_TITLE = Pt(44)
FONT_SUBTITLE = Pt(28)
FONT_BODY = Pt(22)
FONT_SMALL = Pt(16)
FONT_FOOTER = Pt(11)


class PPTXGenerator:
    """Generischer PowerPoint-Generator für Spanisch-Lektionen"""

    def __init__(self, lesson_data, base_dir=None):
        """
        lesson_data: dict mit:
            - id: z.B. '01a'
            - title: Haupttitel
            - subtitle: Untertitel
            - sequenz: Sequenznummer (1-8)
            - block: 'a', 'b', 'c', 'd'
            - slides: Liste von Folien-Definitionen
        base_dir: Basisverzeichnis für relative Pfade (z.B. Bilder)
        """
        self.lesson = lesson_data
        self.base_dir = base_dir or os.getcwd()
        self.prs = None

    def resolve_path(self, path):
        """Löst relativen Pfad relativ zum base_dir auf"""
        if os.path.isabs(path):
            return path
        return os.path.join(self.base_dir, path)

    def create_presentation(self):
        """Erstellt leere 16:9 Präsentation"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        return self.prs

    def add_slide(self):
        """Fügt leere Folie hinzu"""
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_WEISS
        return slide

    # ========================================================================
    # TEXT-HELFER
    # ========================================================================

    def add_title(self, slide, text, top=MARGIN_TOP, size=FONT_TITLE, color=C_SPANISCH):
        """Haupttitel"""
        shape = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Cm(1.8))
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = True
        run.font.color.rgb = color
        return shape

    def add_subtitle(self, slide, text, top=Cm(3.2), size=FONT_SUBTITLE, color=C_TEXT):
        """Untertitel"""
        shape = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Cm(1.5))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        return shape

    def add_text(self, slide, text, left, top, width, height,
                 size=FONT_BODY, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT):
        """Fließtext an beliebiger Position"""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    def add_multiline_text(self, slide, lines, left, top, width,
                           size=FONT_BODY, color=C_TEXT, line_height=1.5):
        """Mehrzeiliger Text"""
        for i, line in enumerate(lines):
            y = top + Cm(line_height * i)
            self.add_text(slide, line, left, y, width, Cm(line_height),
                         size=size, color=color)

    # ========================================================================
    # TABELLEN
    # ========================================================================

    def add_table(self, slide, data, headers, top, col_widths=None):
        """Formatierte Tabelle mit Header"""
        rows = len(data) + 1
        cols = len(headers)

        if col_widths is None:
            col_widths = [Cm(10)] * cols

        total_width = sum(col_widths)
        left = (SLIDE_WIDTH - total_width) / 2
        height = Cm(1.1 * rows)

        table_shape = slide.shapes.add_table(rows, cols, left, top, total_width, height)
        table = table_shape.table

        # Entferne tableStyleId um Repair-Warnung zu vermeiden
        try:
            tbl = table_shape._element.graphic.graphicData.tbl
            tblPr = tbl.tblPr
            if tblPr is not None:
                # Entferne tableStyleId Element
                for child in list(tblPr):
                    if 'tableStyleId' in child.tag:
                        tblPr.remove(child)
        except Exception:
            pass  # Ignorieren wenn es nicht funktioniert

        # Spaltenbreiten
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # Header
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_SPANISCH
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = h
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = C_WEISS
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Daten
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                tf = cell.text_frame
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                run = tf.paragraphs[0].add_run()
                run.text = str(cell_text)
                run.font.size = Pt(20)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Erste Spalte: Spanisch in Rot
                if col_idx == 0:
                    run.font.color.rgb = C_SPANISCH
                    run.font.bold = True
                else:
                    run.font.color.rgb = C_TEXT

                # Zebra-Streifen
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_HELLGRAU

        return table

    # ========================================================================
    # BILDER
    # ========================================================================

    def add_image(self, slide, image_path, left, top, width=None, height=None):
        """Fügt ein Bild zur Folie hinzu"""
        if not os.path.exists(image_path):
            # Fallback: Platzhalter-Box
            self.add_text(slide, f"[Bild: {os.path.basename(image_path)}]",
                         left, top, Cm(10), Cm(5), size=Pt(16), color=C_GRAU)
            return None

        if width and height:
            pic = slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            pic = slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            pic = slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            pic = slide.shapes.add_picture(image_path, left, top)
        return pic

    def create_map_slide(self, title, subtitle, image_path):
        """Folie mit Weltkarte"""
        slide = self.add_slide()

        self.add_title(slide, title)
        if subtitle:
            self.add_subtitle(slide, subtitle)

        # Relativen Pfad auflösen
        resolved_path = self.resolve_path(image_path)

        # Bild zentriert einfügen
        if os.path.exists(resolved_path):
            # Bildgröße: ca. 26cm breit, zentriert
            img_width = Cm(26)
            img_left = (SLIDE_WIDTH - img_width) / 2
            self.add_image(slide, resolved_path, img_left, Cm(5), width=img_width)
        else:
            self.add_text(slide, f"[Karte nicht gefunden: {image_path}]",
                         Cm(5), Cm(8), Cm(24), Cm(3), size=Pt(20), color=C_GRAU,
                         align=PP_ALIGN.CENTER)

        self.add_footer(slide)
        return slide

    # ========================================================================
    # SPEZIELLE BOXEN
    # ========================================================================

    def add_info_box(self, slide, title, content, top, icon="💡"):
        """Info-Box mit orangem Rahmen"""
        left = Cm(2)
        width = Cm(29.8)
        height = Cm(3.5)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xff, 0xf8, 0xf0)
        rect.line.color.rgb = C_ORANGE
        rect.line.width = Pt(2)

        self.add_text(slide, f"{icon} {title}", left + Cm(0.5), top + Cm(0.3),
                      width - Cm(1), Cm(1), size=Pt(18), color=C_ORANGE, bold=True)
        self.add_text(slide, content, left + Cm(0.5), top + Cm(1.3),
                      width - Cm(1), height - Cm(1.5), size=Pt(18), color=C_TEXT)

    def add_ab_reference(self, slide, text, top):
        """AB-Referenz-Box (karminrot)"""
        left = Cm(2)
        width = Cm(29.8)
        height = Cm(1.8)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xfd, 0xf2, 0xf4)
        rect.line.color.rgb = C_SPANISCH
        rect.line.width = Pt(1.5)

        self.add_text(slide, f"📝 {text}", left + Cm(0.5), top + Cm(0.35),
                      width - Cm(1), Cm(1.2), size=Pt(16), color=C_SPANISCH, bold=True)

    def add_grammar_box(self, slide, title, content, top):
        """Grammatik-Box (grün)"""
        left = Cm(2)
        width = Cm(29.8)
        height = Cm(4)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xf0, 0xf8, 0xf0)
        rect.line.color.rgb = C_GRUEN
        rect.line.width = Pt(2)

        self.add_text(slide, f"📝 {title}", left + Cm(0.5), top + Cm(0.3),
                      width - Cm(1), Cm(1), size=Pt(20), color=C_GRUEN, bold=True)
        self.add_text(slide, content, left + Cm(0.5), top + Cm(1.5),
                      width - Cm(1), height - Cm(2), size=Pt(18), color=C_TEXT)

    def add_footer(self, slide):
        """Standard-Fußzeile"""
        lesson_id = self.lesson.get('id', '??')
        left_text = f"Spanisch Klasse 7 · LP-{lesson_id}"
        right_text = "Greenhouse Schools"

        self.add_text(slide, left_text, MARGIN_LEFT, FOOTER_TOP, Cm(12), Cm(0.8),
                      size=FONT_FOOTER, color=C_GRAU)
        self.add_text(slide, right_text, Cm(20), FOOTER_TOP, Cm(12), Cm(0.8),
                      size=FONT_FOOTER, color=C_GRAU, align=PP_ALIGN.RIGHT)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_LEFT, FOOTER_TOP - Cm(0.15),
                                       CONTENT_WIDTH, Cm(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = C_HELLGRAU
        line.line.fill.background()

    # ========================================================================
    # STANDARD-FOLIEN
    # ========================================================================

    def create_title_slide(self, title, subtitle, meta=None, flags=None):
        """Titelfolie"""
        slide = self.add_slide()

        self.add_text(slide, title, Cm(2), Cm(5), Cm(30), Cm(3),
                      size=Pt(72), color=C_SPANISCH, bold=True, align=PP_ALIGN.CENTER)

        self.add_text(slide, subtitle, Cm(2), Cm(9.5), Cm(30), Cm(1.5),
                      size=Pt(32), color=C_TEXT, align=PP_ALIGN.CENTER)

        if meta:
            self.add_text(slide, meta, Cm(2), Cm(11.5), Cm(30), Cm(1),
                          size=Pt(20), color=C_GRAU, align=PP_ALIGN.CENTER)

        if flags:
            self.add_text(slide, flags, Cm(2), Cm(14), Cm(30), Cm(1.5),
                          size=Pt(36), color=C_TEXT, align=PP_ALIGN.CENTER)

        self.add_footer(slide)
        return slide

    def create_vocab_slide(self, title, subtitle, vocab_data, ab_ref=None):
        """Vokabel-Tabelle"""
        slide = self.add_slide()

        self.add_title(slide, title)
        if subtitle:
            self.add_subtitle(slide, subtitle)

        # Bestimme Headers basierend auf Daten
        if vocab_data and len(vocab_data[0]) == 3:
            headers = ["Spanisch", "Deutsch", "Hinweis"]
            widths = [Cm(11), Cm(9), Cm(8)]
        else:
            headers = ["Spanisch", "Deutsch"]
            widths = [Cm(12), Cm(14)]

        self.add_table(slide, vocab_data, headers, Cm(5), col_widths=widths)

        if ab_ref:
            self.add_ab_reference(slide, ab_ref, Cm(14))

        self.add_footer(slide)
        return slide

    def create_exercise_slide(self, title, subtitle, exercises, options=None, ab_ref=None):
        """Übungsfolie"""
        slide = self.add_slide()

        self.add_title(slide, f"✏️ {title}")
        if subtitle:
            self.add_subtitle(slide, subtitle)

        for i, (left_text, right_text) in enumerate(exercises):
            top = Cm(5.5 + i * 2.2)
            self.add_text(slide, left_text, Cm(4), top, Cm(12), Cm(1.2),
                          size=Pt(26), color=C_SPANISCH, bold=True)
            self.add_text(slide, right_text, Cm(17), top, Cm(12), Cm(1.2),
                          size=Pt(26), color=C_GRAU)

        if options:
            self.add_text(slide, f"Optionen: {options}", Cm(4), Cm(14.5), Cm(26), Cm(1),
                          size=Pt(18), color=C_GRAU)

        if ab_ref:
            self.add_ab_reference(slide, ab_ref, Cm(16))

        self.add_footer(slide)
        return slide

    def create_solution_slide(self, title, solutions):
        """Lösungsfolie"""
        slide = self.add_slide()

        self.add_title(slide, f"✅ {title}")

        for i, (left_text, right_text) in enumerate(solutions):
            top = Cm(4.5 + i * 2.5)
            self.add_text(slide, left_text, Cm(4), top, Cm(12), Cm(1.2),
                          size=Pt(28), color=C_SPANISCH, bold=True)
            self.add_text(slide, right_text, Cm(16), top, Cm(14), Cm(1.2),
                          size=Pt(28), color=C_GRUEN, bold=True)

        self.add_footer(slide)
        return slide

    def create_grammar_slide(self, title, content_blocks, info_box=None):
        """Grammatik-Folie"""
        slide = self.add_slide()

        self.add_title(slide, f"📝 {title}")

        y = Cm(4)
        for block in content_blocks:
            if isinstance(block, tuple) and len(block) == 2:
                # (spanisch, erklärung)
                self.add_text(slide, block[0], Cm(3), y, Cm(14), Cm(1.5),
                              size=Pt(36), color=C_SPANISCH, bold=True)
                self.add_text(slide, block[1], Cm(18), y + Cm(0.3), Cm(14), Cm(1),
                              size=Pt(22), color=C_GRAU)
                y += Cm(2.5)
            else:
                # Einfacher Text
                self.add_text(slide, str(block), Cm(3), y, Cm(28), Cm(1.2),
                              size=Pt(24), color=C_TEXT)
                y += Cm(2)

        if info_box:
            self.add_info_box(slide, info_box[0], info_box[1], Cm(12))

        self.add_footer(slide)
        return slide

    def create_summary_slide(self, title, columns):
        """Zusammenfassungs-Folie mit 2 Spalten"""
        slide = self.add_slide()

        self.add_title(slide, f"📋 {title}")

        # Linke Spalte
        if len(columns) >= 1:
            col1 = columns[0]
            self.add_text(slide, col1['title'], Cm(2.5), Cm(4), Cm(14), Cm(1.2),
                          size=Pt(24), color=C_SPANISCH, bold=True)
            for i, item in enumerate(col1['items']):
                self.add_text(slide, f"•  {item}", Cm(3), Cm(5.5 + i * 1.6), Cm(12), Cm(1.2),
                              size=Pt(22), color=C_TEXT)

        # Rechte Spalte
        if len(columns) >= 2:
            col2 = columns[1]
            self.add_text(slide, col2['title'], Cm(17.5), Cm(4), Cm(14), Cm(1.2),
                          size=Pt(24), color=C_SPANISCH, bold=True)
            for i, item in enumerate(col2['items']):
                self.add_text(slide, f"•  {item}", Cm(18), Cm(5.5 + i * 1.6), Cm(12), Cm(1.2),
                              size=Pt(22), color=C_TEXT)

        self.add_footer(slide)
        return slide

    def create_homework_slide(self, title, tasks, ab_ref=None, closing=None):
        """Hausaufgaben-Folie"""
        slide = self.add_slide()

        self.add_title(slide, f"📚 {title}")

        for i, task in enumerate(tasks):
            top = Cm(4.5 + i * 2.3)
            self.add_text(slide, "✓", Cm(3), top, Cm(1.5), Cm(1.2),
                          size=Pt(28), color=C_GRUEN, bold=True)
            self.add_text(slide, task, Cm(5), top, Cm(26), Cm(1.2),
                          size=Pt(24), color=C_TEXT)

        if ab_ref:
            self.add_ab_reference(slide, ab_ref, Cm(14))

        if closing:
            self.add_text(slide, closing, Cm(2), Cm(16.5), Cm(30), Cm(1.5),
                          size=Pt(36), color=C_SPANISCH, bold=True, align=PP_ALIGN.CENTER)

        self.add_footer(slide)
        return slide

    def create_content_slide(self, title, content, ab_ref=None):
        """Generische Inhaltsfolie"""
        slide = self.add_slide()

        self.add_title(slide, title)

        y = Cm(4)
        for item in content:
            if isinstance(item, dict):
                if item.get('type') == 'text':
                    self.add_text(slide, item['text'],
                                  Cm(item.get('left', 3)), Cm(item.get('top', y)),
                                  Cm(item.get('width', 28)), Cm(item.get('height', 1.2)),
                                  size=Pt(item.get('size', 24)),
                                  color=item.get('color', C_TEXT),
                                  bold=item.get('bold', False))
                elif item.get('type') == 'table':
                    self.add_table(slide, item['data'], item['headers'],
                                   Cm(item.get('top', 5)),
                                   col_widths=[Cm(w) for w in item.get('widths', [10, 10])])
                elif item.get('type') == 'info_box':
                    self.add_info_box(slide, item['title'], item['content'],
                                      Cm(item.get('top', 12)))
                elif item.get('type') == 'grammar_box':
                    self.add_grammar_box(slide, item['title'], item['content'],
                                         Cm(item.get('top', 8)))
            else:
                self.add_text(slide, str(item), Cm(3), y, Cm(28), Cm(1.2),
                              size=Pt(24), color=C_TEXT)
                y += Cm(2)

        if ab_ref:
            self.add_ab_reference(slide, ab_ref, Cm(16))

        self.add_footer(slide)
        return slide

    # ========================================================================
    # GENERIERUNG
    # ========================================================================

    def generate(self, output_path):
        """Generiert Präsentation aus Lesson-Daten"""
        self.create_presentation()

        slides_data = self.lesson.get('slides', [])

        for slide_def in slides_data:
            slide_type = slide_def.get('type', 'content')

            if slide_type == 'title':
                self.create_title_slide(
                    slide_def.get('title', ''),
                    slide_def.get('subtitle', ''),
                    slide_def.get('meta'),
                    slide_def.get('flags')
                )

            elif slide_type == 'vocab':
                self.create_vocab_slide(
                    slide_def.get('title', 'Vokabeln'),
                    slide_def.get('subtitle'),
                    slide_def.get('vocab', []),
                    slide_def.get('ab_ref')
                )

            elif slide_type == 'exercise':
                self.create_exercise_slide(
                    slide_def.get('title', 'Übung'),
                    slide_def.get('subtitle'),
                    slide_def.get('exercises', []),
                    slide_def.get('options'),
                    slide_def.get('ab_ref')
                )

            elif slide_type == 'solution':
                self.create_solution_slide(
                    slide_def.get('title', 'Lösung'),
                    slide_def.get('solutions', [])
                )

            elif slide_type == 'grammar':
                self.create_grammar_slide(
                    slide_def.get('title', 'Grammatik'),
                    slide_def.get('content', []),
                    slide_def.get('info_box')
                )

            elif slide_type == 'summary':
                self.create_summary_slide(
                    slide_def.get('title', 'Zusammenfassung'),
                    slide_def.get('columns', [])
                )

            elif slide_type == 'homework':
                self.create_homework_slide(
                    slide_def.get('title', 'Hausaufgabe'),
                    slide_def.get('tasks', []),
                    slide_def.get('ab_ref'),
                    slide_def.get('closing')
                )

            elif slide_type == 'map':
                self.create_map_slide(
                    slide_def.get('title', 'Weltkarte'),
                    slide_def.get('subtitle'),
                    slide_def.get('image_path', '')
                )

            elif slide_type == 'content':
                self.create_content_slide(
                    slide_def.get('title', ''),
                    slide_def.get('content', []),
                    slide_def.get('ab_ref')
                )

        self.prs.save(output_path)

        # Post-Processing: Entferne def-Attribut aus tableStyles.xml
        self._fix_table_styles(output_path)

        return len(self.prs.slides)

    def _fix_table_styles(self, pptx_path):
        """Entfernt alle problematischen Tabellen-Style-Referenzen aus der PPTX"""
        import zipfile
        import tempfile
        import shutil
        import re

        try:
            temp_dir = tempfile.mkdtemp()
            temp_pptx = os.path.join(temp_dir, 'temp.pptx')

            with zipfile.ZipFile(pptx_path, 'r') as zip_in:
                with zipfile.ZipFile(temp_pptx, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for item in zip_in.infolist():
                        data = zip_in.read(item.filename)

                        # tableStyles.xml: Entferne def-Attribut
                        if item.filename == 'ppt/tableStyles.xml':
                            data_str = data.decode('utf-8')
                            data_str = re.sub(r'\s*def="[^"]*"', '', data_str)
                            data = data_str.encode('utf-8')

                        # Slide-XML: Entferne tableStyleId-Elemente
                        elif item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                            data_str = data.decode('utf-8')
                            # Entferne <a:tableStyleId>...</a:tableStyleId>
                            data_str = re.sub(r'<a:tableStyleId>[^<]*</a:tableStyleId>', '', data_str)
                            data = data_str.encode('utf-8')

                        zip_out.writestr(item, data)

            shutil.move(temp_pptx, pptx_path)
            shutil.rmtree(temp_dir)
        except Exception as e:
            pass  # Bei Fehler Original behalten


# ============================================================================
# BEISPIEL-VERWENDUNG
# ============================================================================

if __name__ == "__main__":
    # Test mit einfachem Beispiel
    test_lesson = {
        'id': 'TEST',
        'title': 'Test-Präsentation',
        'subtitle': 'Zum Testen des Generators',
        'sequenz': 0,
        'block': 'x',
        'slides': [
            {
                'type': 'title',
                'title': '¡Hola!',
                'subtitle': 'Test-Präsentation',
                'meta': 'Sequenz 0 · Test',
                'flags': '🇪🇸'
            },
            {
                'type': 'vocab',
                'title': '📚 Testvokabeln',
                'subtitle': 'Einige Beispielwörter',
                'vocab': [
                    ('Hola', 'Hallo'),
                    ('Adiós', 'Tschüss'),
                ],
                'ab_ref': 'AB-TEST: Kasten 1'
            },
            {
                'type': 'homework',
                'title': 'Hausaufgabe',
                'tasks': ['Vokabeln lernen', 'AB ausfüllen'],
                'closing': '¡Hasta luego!'
            }
        ]
    }

    gen = PPTXGenerator(test_lesson)
    count = gen.generate('test_pptx.pptx')
    print(f"Test-Präsentation erstellt mit {count} Folien")
